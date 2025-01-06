import PyPDF2
import pandas as pd
from docx import Document
from pptx import Presentation


class FileReader:
    def __init__(self):
        pass

    def check_exists(self, file_path):
        try:
            with open(file_path, 'rb') as file:
                return True
        except FileNotFoundError:
            return

    def read_pdf(self, file_path):
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text()
        return text
    
    def read_text(self, file_path):
        with open(file_path, 'rb') as file:
            return file.read()
        
    def read_sheet(self, file_path):
        df = ""
        if file_path.endswith(".csv"):
            df = pd.read_csv(file_path)
        elif file_path.endswith(".xlsx"):
            df = pd.read_excel(file_path)
        return df.to_string(index=False)
    
    def read_docx(self, file_path):
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text
        return '\n'.join(text)
    
    def read_pptx(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        return '\n'.join(text)

    def reader(self, file_path):
        if file_path.endswith(".pdf"):
            return self.read_pdf(file_path)
        elif file_path.endswith(".csv") or file_path.endswith(".xlsx"):
            return self.read_sheet(file_path)
        elif file_path.endswith(".docx") or file_path.endswith(".doc"):
            return self.read_docx(file_path)
        elif file_path.endswith(".pptx") or file_path.endswith(".ppt"):
            return self.read_pptx(file_path)
        else:
            return self.read_text(file_path)
